import os
import glob
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from PIL import Image

def points_to_pixels(points):
    return points * 96 / 72

def autofit_text_shape(shape):
    text_frame = shape.text_frame
    text_frame.word_wrap = False
    text_frame.auto_size = True

# New function to get image files from folder
def get_image_files(image_folder):
    extensions = ["*.jpg", "*.jpeg", "*.png"]
    image_files = []

    for ext in extensions:
        image_files.extend(glob.glob(os.path.join(image_folder, ext)))

    return image_files

# New function to create the title slide
def create_title_slide(pptx, title_input, subtitle_input):
    title_slide_layout = pptx.slide_layouts[0]
    slide = pptx.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = str(title_input)
    subtitle.text = str(subtitle_input)

def create_image_slide(pptx, image_file):
    slide_layout = pptx.slide_layouts[6]
    slide = pptx.slides.add_slide(slide_layout)

    with Image.open(image_file) as img:
        width, height = img.size
        aspect_ratio = float(width) / float(height)

    slide_width = pptx.slide_width
    slide_height = pptx.slide_height
    slide_width_px = points_to_pixels(slide_width.pt)
    slide_height_px = points_to_pixels(slide_height.pt)

    if aspect_ratio > 1:
        width = slide_width_px
        height = int(slide_width_px / aspect_ratio)
    else:
        height = slide_height_px
        width = int(slide_height_px * aspect_ratio)

    left = (slide_width - Inches(width / 96)) / 2
    top = (slide_height - Inches(height / 96)) / 2
    pic = slide.shapes.add_picture(image_file, left, top, width=Inches(width / 96), height=Inches(height / 96))

    # Add the title text box to the slide
    image_name = os.path.basename(image_file)
    char_count = len(image_name)
    fixed_width_font_size = Inches(0.2)  # Adjust the font size here
    char_width_inches = 0.13  # Approximate width of a character in inches for a fixed-width font

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(char_count * char_width_inches), Inches(0.35)
    )
    title = title_box.text_frame
    title.text = image_name

    # Set the font properties
    title.paragraphs[0].runs[0].font.size = fixed_width_font_size
    title.paragraphs[0].runs[0].font.name = "Courier New"  # Fixed-width font

    autofit_text_shape(title_box)

    # Set the background color of the text box to white
    fill = title_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

def create_presentation(image_folder, output_file):
    pptx = Presentation()
    image_folder = os.path.abspath(image_folder)
    image_files = get_image_files(image_folder)

    title_input = input("Enter a title: ")
    subtitle_input = input("Enter a subtitle: ")

    create_title_slide(pptx, title_input, subtitle_input)

    for image_file in image_files:
        create_image_slide(pptx, image_file)

    pptx.save(output_file)

if __name__ == "__main__":
    image_folder = os.path.dirname(os.path.realpath(__file__))
    output_file = os.path.join(image_folder, "presentation.pptx")
    create_presentation(image_folder, output_file)
    os.startfile(output_file)
